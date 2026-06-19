"""
Fill ISRO BAH 2026 Idea Submission Template with LunarIceNet content.
Embeds visual assets (diagrams, charts) + actual pipeline output PNGs.

Run: cd presentation && python fill_presentation.py
Output: LunarIceNet_BAH2026_PS8.pptx
"""
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path

SRC = "[Pub] ISRO BAH 2026 _ Idea Submission Template.pptx"
OUT = "LunarIceNet_BAH2026_PS8.pptx"
ASSETS = Path("assets")
OUTPUTS = Path("..") / "outputs"

# ── helpers ────────────────────────────────────────────────────────────────────

def set_textbox(slide, shape_id, lines, title_line=None, font_size_pt=None):
    """Replace text in a shape (by shape_id) while preserving background/layout."""
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            tf = shape.text_frame
            tf.word_wrap = True

            sample_font_name = "Google Sans"
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        sample_font_name = run.font.name
                    break
                break

            use_size = Pt(font_size_pt) if font_size_pt else None

            txBody = tf._txBody
            for p in txBody.findall(qn('a:p')):
                txBody.remove(p)

            def add_para(text, bold=False, size_pt=None, color_rgb=None):
                p_elem = etree.SubElement(txBody, qn('a:p'))
                r_elem = etree.SubElement(p_elem, qn('a:r'))
                rPr = etree.SubElement(r_elem, qn('a:rPr'), lang='en-US', dirty='0')
                sz = size_pt if size_pt else 13
                rPr.set('sz', str(int(sz * 100)))
                if bold:
                    rPr.set('b', '1')
                if color_rgb:
                    solidFill = etree.SubElement(rPr, qn('a:solidFill'))
                    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
                    srgbClr.set('val', color_rgb)
                latin = etree.SubElement(rPr, qn('a:latin'))
                latin.set('typeface', sample_font_name)
                t_elem = etree.SubElement(r_elem, qn('a:t'))
                t_elem.text = text

            if title_line:
                add_para(title_line, bold=True, size_pt=22, color_rgb='FFFFFF')

            for line in lines:
                if line == "":
                    add_para("", size_pt=8)
                elif line.startswith("##"):
                    add_para(line[2:].strip(), bold=True, size_pt=16, color_rgb='FFD700')
                elif line.startswith("**") and line.endswith("**"):
                    add_para(line[2:-2], bold=True, size_pt=13, color_rgb='00E5FF')
                elif line.startswith("\u2022 "):
                    add_para(line, size_pt=12)
                else:
                    add_para(line, size_pt=12)
            return
    print(f"WARNING: shape_id {shape_id} not found on slide")


def add_image(slide, img_path, left, top, width=None, height=None):
    """Add an image to a slide."""
    p = Path(img_path)
    if not p.exists():
        print(f"WARNING: image not found: {p}")
        return
    kwargs = {"image_file": str(p), "left": Inches(left), "top": Inches(top)}
    if width:
        kwargs["width"] = Inches(width)
    if height:
        kwargs["height"] = Inches(height)
    slide.shapes.add_picture(**kwargs)


# ── content ────────────────────────────────────────────────────────────────────

prs = Presentation(SRC)

# ── SLIDE 1: Cover ─────────────────────────────────────────────────────────────
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.shape_id == 55:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = "Team Name : LunarIceNet"
    elif shape.shape_id == 56:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = "Problem Statement : PS-8 \u2014 Subsurface Ice Detection in Lunar South Polar Regions"
    elif shape.shape_id == 57:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = "Team Leader Name : Priyanshu Doshi"

# ── SLIDE 2: Team Members ──────────────────────────────────────────────────────
slide2 = prs.slides[1]
for shape in slide2.shapes:
    if shape.shape_type == 19:  # TABLE
        tbl = shape.table

        def fill_cell(row, col, title, name, college="\u2014"):
            cell = tbl.cell(row, col)
            tf = cell.text_frame
            tf.word_wrap = True
            txBody = tf._txBody
            for p in txBody.findall(qn('a:p')):
                txBody.remove(p)
            for line, bold in [(title, True), ("", False), (f"Name: {name}", False), (f"College: {college}", False)]:
                p_elem = etree.SubElement(txBody, qn('a:p'))
                r_elem = etree.SubElement(p_elem, qn('a:r'))
                rPr = etree.SubElement(r_elem, qn('a:rPr'), lang='en-US', dirty='0')
                rPr.set('sz', '1400' if bold else '1200')
                if bold:
                    rPr.set('b', '1')
                latin = etree.SubElement(rPr, qn('a:latin'))
                latin.set('typeface', 'Google Sans')
                t_elem = etree.SubElement(r_elem, qn('a:t'))
                t_elem.text = line

        fill_cell(0, 0, "Team Leader:",   "Priyanshu Doshi")
        fill_cell(0, 1, "Team Member 1:", "Shub Patel")
        fill_cell(1, 0, "Team Member 2:", "Meer Patel")
        fill_cell(1, 1, "Team Member 3:", "\u2014")

# ── SLIDE 3: Opportunity / USP ────────────────────────────────────────────────
set_textbox(prs.slides[2], 70, [
    "## THE PROBLEM",
    "",
    "\u2022 Chandrayaan-2 DFSAR captured the first full-polarimetric L-band SAR of the lunar south pole",
    "\u2022 140M+ radar pixels \u2014 the largest lunar SAR dataset ever acquired",
    "\u2022 Classical CPR > 1.0 thresholding gives high false positives and zero depth information",
    "\u2022 No end-to-end pipeline exists: raw DFSAR \u2192 ice map \u2192 landing site \u2192 rover path \u2192 ice volume",
    "",
    "## OUR SOLUTION: LunarIceNet",
    "",
    "\u2022 First deep learning system trained directly on real Chandrayaan-2 DFSAR Level 3C mosaics",
    "\u2022 Physics-informed loss penalizes predictions violating temperature/latitude constraints",
    "\u2022 Cross-attention fuses radar features with physical context (lat, PSR membership)",
    "\u2022 Multi-task output: ice probability + penetration depth + confidence in ONE forward pass",
    "\u2022 Single-command pipeline: raw data \u2192 complete mission report in 15 minutes",
    "",
    "## WHY IT WINS",
    "",
    "\u2022 Best F1 = 0.938 on real DFSAR data (west model) \u2014 not synthetic benchmarks",
    "\u2022 93.3% real terrain coverage from NASA LOLA DEM (not made-up slopes)",
    "\u2022 Monte Carlo uncertainty quantification \u2014 honest error bars, not fake precision",
    "\u2022 Validated against PRL 2026 (Sinha et al.) \u2014 replicates published science",
])
# Add key metrics infographic at bottom
add_image(prs.slides[2], ASSETS / "key_metrics.png", left=0.3, top=5.5, width=9.4)

# ── SLIDE 4: Features ─────────────────────────────────────────────────────────
set_textbox(prs.slides[3], 76, [
    "## 5 CORE CAPABILITIES",
    "",
    "**1. Subsurface Ice Detection (LunarIceNet CNN \u2014 12.4M params)**",
    "\u2022 Input: CPR, SERD, T-Ratio from DFSAR Level 3C",
    "\u2022 Multi-scale ResNet + Squeeze-Excitation + Cross-Attention fusion",
    "\u2022 Output: per-pixel ice probability, depth (m), confidence maps",
    "",
    "**2. Real Terrain Integration (NASA LOLA DEM)**",
    "\u2022 93.3% coverage at 40m/pixel \u2014 real slopes, not proxies",
    "",
    "**3. Landing Site Selection (Multi-criteria scoring)**",
    "\u2022 35% ice + 20% slope + 15% access + 15% illumination + 15% confidence",
    "\u2022 4000+ candidates evaluated \u2192 top 10 ranked",
    "",
    "**4. Rover Traverse Planning (A* pathfinding)**",
    "\u2022 8-connected grid, cost = distance + slope_pen + solar_pen \u2212 ice_reward",
    "\u2022 Hard limit: slope > 25\u00b0 impassable, 8 km range budget",
    "",
    "**5. Ice Volume Estimation (Lichtenecker + Monte Carlo)**",
    "\u2022 CPR \u2192 ice fraction \u2192 penetration depth \u2192 volume per pixel",
    "\u2022 1000-sample Monte Carlo for 90% confidence interval",
    "\u2022 Per-crater breakdown across 16 known DPSRs",
])
# Add results comparison chart
add_image(prs.slides[3], ASSETS / "results_comparison.png", left=0.3, top=5.2, width=9.4)

# ── SLIDE 5: Process Flow ─────────────────────────────────────────────────────
set_textbox(prs.slides[4], 82, [
    "## END-TO-END PIPELINE",
    "",
    "\u2022 Single command: python full_pipeline.py --direction west",
    "\u2022 Fully automated: data loading \u2192 inference \u2192 analysis \u2192 visualization \u2192 report",
    "\u2022 Runs on CPU (15 min cached) or GPU (2 min inference)",
])
# Add pipeline flow diagram
add_image(prs.slides[4], ASSETS / "pipeline_flow.png", left=0.2, top=2.8, width=9.6)
# Add actual output: ice analysis map
ice_map = OUTPUTS / "ice_analysis_east.png"
if ice_map.exists():
    add_image(prs.slides[4], ice_map, left=0.5, top=5.2, width=4.0)
# Add polar map
polar_map = OUTPUTS / "polar_ice_map_east.png"
if polar_map.exists():
    add_image(prs.slides[4], polar_map, left=5.0, top=5.2, width=4.5)

# ── SLIDE 6: Output Visualizations ─────────────────────────────────────────
set_textbox(prs.slides[5], 88, [
    "## ACTUAL PIPELINE OUTPUTS",
])
# Embed real output images
traverse = OUTPUTS / "rover_traverse_east.png"
if traverse.exists():
    add_image(prs.slides[5], traverse, left=0.3, top=1.5, width=4.5)
west_ice = OUTPUTS / "west" / "ice_analysis_west.png"
if west_ice.exists():
    add_image(prs.slides[5], west_ice, left=5.0, top=1.5, width=4.5)
# Add LOLA DEM
lola = OUTPUTS / "lola_dem_summary.png"
if lola.exists():
    add_image(prs.slides[5], lola, left=0.3, top=5.0, width=4.5)
west_polar = OUTPUTS / "west" / "polar_ice_map_west.png"
if west_polar.exists():
    add_image(prs.slides[5], west_polar, left=5.0, top=5.0, width=4.5)

# ── SLIDE 7: Architecture ─────────────────────────────────────────────────────
set_textbox(prs.slides[6], 94, [
    "## LUNARICENET ARCHITECTURE",
])
# Embed architecture diagram
add_image(prs.slides[6], ASSETS / "architecture_diagram.png", left=0.2, top=1.3, width=9.6)

# ── SLIDE 8: Technologies ─────────────────────────────────────────────────────
set_textbox(prs.slides[7], 100, [
    "## TECHNOLOGY STACK",
    "",
    "**Deep Learning:**  PyTorch 2.2 + AMP mixed precision + torchmetrics",
    "**Geospatial:**  rasterio (GeoTIFF/UPS), numpy, scipy (slope computation)",
    "**Data Sources:**  ISRO PRADAN (DFSAR L3C) + NASA PDS (LOLA GDR DEMs)",
    "**Planning:**  Custom A* (8-connected grid) + Monte Carlo uncertainty",
    "**Physics:**  Lichtenecker dielectric mixing + CPR \u2192 ice fraction curve",
    "**Visualization:**  matplotlib (polar stereo, multi-panel, dark theme)",
    "**Infrastructure:**  NVIDIA CUDA/cuDNN, Python 3.12, Git/GitHub",
    "",
    "**Data Requirements**",
    "\u2022 DFSAR Level 3C mosaics: ~2 GB (east) + ~3 GB (west) \u2014 free from PRADAN",
    "\u2022 LOLA DEM: ~110 MB (ldem_85s_40m) \u2014 free from NASA PDS",
    "\u2022 Total storage: ~5 GB for data + checkpoints + outputs",
    "",
    "**Hardware Used**",
    "\u2022 Training: NVIDIA GPU (8+ GB VRAM), 32 GB RAM",
    "\u2022 Inference: CPU-only capable (15 min per direction)",
    "\u2022 All open-source \u2014 zero license fees",
])

# ── SLIDE 9: Cost ─────────────────────────────────────────────────────────────
set_textbox(prs.slides[8], 106, [
    "## COST & FEASIBILITY",
    "",
    "**Development Cost**",
    "\u2022 Model training (30 epochs \u00d7 2 directions): ~8 hrs GPU \u2248 \u20b9120 cloud or \u20b950 own GPU",
    "\u2022 All data sources are FREE (ISRO PRADAN + NASA PDS)",
    "\u2022 All software is open-source (PyTorch, rasterio, numpy, scipy)",
    "",
    "**Deployment Cost**",
    "\u2022 Inference on new DFSAR mosaic: ~15 min CPU or ~2 min GPU",
    "\u2022 No internet required after data download \u2014 fully offline capable",
    "\u2022 Single Python script \u2014 no cloud infrastructure needed",
    "",
    "**Chandrayaan-4 Integration Path**",
    "\u2022 Drop-in for ISRO ground segment: GeoTIFF + mission report output",
    "\u2022 Retrain on new instrument data with same pipeline",
    "\u2022 Real-time onboard inference possible with model quantization (INT8)",
    "",
    "**Total cost to reproduce: < \u20b9500**",
    "(all data free, all software open-source, runs on consumer hardware)",
    "",
    "## FUTURE WORK",
    "",
    "\u2022 Multi-resolution depth profiling (L-band deep + S-band shallow)",
    "\u2022 Interactive Streamlit dashboard for mission planning",
    "\u2022 GeoTIFF export for QGIS/ArcGIS integration",
    "\u2022 Ensemble east+west predictions for robust ice mapping",
    "\u2022 Transfer learning to Chandrayaan-4 instruments",
])

# ── SAVE ──────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
print("Done! Open in PowerPoint to review.")
